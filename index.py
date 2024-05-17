import pptx

import work_with_exel
import work_with_pptx
import work_with_json

import pandas as pd


from pptx import Presentation
from pptx.util import Inches, Pt


def main() -> None:
    # Получаю данные из exel и создаю презентацию в pptx
    # в будущем создам функции для работы с pptx
    json_data = work_with_json.read_json_file("./files.json")

    name_exel_file = json_data["exel_in_file"]
    xl_file = pd.ExcelFile(name_exel_file)
    name_sheets = xl_file.sheet_names

    prs = Presentation()

    ind = 0
    tables = []
    for name in name_sheets:
        if ind % 2 == 0:
            ind += 1
            continue

        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = f"Adding a Table {ind}"

        stat = work_with_exel.get_data_from_sheet(name_exel_file, name)
        num_rows, num_columns = stat.shape

        rows = num_rows
        columns = num_columns
        left = top = Inches(1.5)
        width = Inches(6.0)
        height = Inches(0.4)
        table = shapes.add_table(rows, columns, left, top, width, height).table

        for row in stat.itertuples():
            for index in range(stat.columns.size):
                value = str(stat.iloc[row.Index, index])
                table.cell(int(row.Index), int(index)).text = value
                table.cell(int(row.Index), int(index)).text_frame.paragraphs[0].runs[0].font.size = Pt(10)

        tables.append(stat)
        ind += 1

    prs.save('test.pptx')
    print(tables)


if __name__ == "__main__":
    main()

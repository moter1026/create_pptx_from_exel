#TODO: файл пока сложно обработать, оставлю напоследок, если время останется. @nick-vivo
# plak plak ;-(
import pandas as pd

from pptx import Presentation
from pptx.util import Inches


class Present:

    table_top = Inches(0.35)
    table_width = Inches(4.0)
    table_height = Inches(0.8)
    table_lefts = [Inches(2.95), Inches(3) + table_width, Inches(3.05) + 2 * table_width]


    def __init__(self, name_of_file: str, template_file: str):

        self.name_of_file = name_of_file
        self.prs = Presentation(template_file)


    def add_slide(self, data: dict) -> None:
        
        images = data["images"]
        tables = data["tables"]
        texts = data["texts"]
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])

        for table in tables:
            num_rows, num_columns = table.shape
            slide.shapes.add_table(num_rows, num_columns, Inches(2), Inches(2), Inches(4), Inches(1.5))


    def save_titul_slide(self, group: str) -> None:
        
        """
        Меняет название группы на титульном слайде шаблона
        """
        
        shapes = self.prs.slides[0].shapes
        for shape in shapes:
            # Check if the shape is a text box
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if '{group}' in text_frame.text:
                    text_frame.text = text_frame.text.replace("{group}", group)


    def add_table_to_slide(self, data: pd.DataFrame, slide_index: int) -> None:
        
        """
        Добавляет подготовленную таблицу на слайд шаблонной презы
        :param data: подготволенная таблица из трех столбцов
        :param slide_index: номер слайда, считая от нуля
        """

        shapes = self.prs.slides[slide_index].shapes
        rows, columns = data.shape

        # определяем количество таблиц на листе
        count_of_tables = rows // 16
        if rows - count_of_tables * 16 > 0:
            count_of_tables += 1

        tables = []
        names = data.columns
        for table_i in range(0, count_of_tables):
            tables.append(shapes.add_table(17, columns, self.table_lefts[table_i], self.table_top, self.table_width,
                                           self.table_height).table)
            for i in range(0, len(names)):
                tables[table_i].cell(0, i).text = names[i]
        for row in data.itertuples():
            for index in range(data.columns.size):
                value = str(data.iloc[row.Index, index])
                tables[row.Index // 16].cell(int(row.Index % 16 + 1), int(index)).text = value

    def add_image_to_slide(self, img_path: str, slide_index: int) -> None:

        """
        Добавляет подготовленные изображение на слайд шаблонной презы
        :param img_path:
        :param slide_index: номер слайда, куда необходимо вставить
        :return: None
        """

        shapes = self.prs.slides[slide_index].shapes
        left = Inches(3.8)  # Расположение по горизонтали
        top = Inches(0.8)  # Расположение по вертикали
        width, height = (Inches(8.5), Inches(6.39))

        shapes.add_picture(img_path, left, top, width=width, height=height)

    def add_mini_table_to_slide(self, data: dict, slide_index: int) -> None:

        top = Inches(4.58)
        left = Inches(2.55)
        width = Inches(5)
        height = Inches(1.2)
        shapes = self.prs.slides[slide_index].shapes
        table = shapes.add_table(2, len(data.keys()), left, top, width,
                                 height).table
        keys = list(data.keys())
        for i in range(0, len(keys)):
            table.cell(0, i).text = keys[i]
            table.cell(1, i).text = data[keys[i]]

    def add_last_tables(self, data: pd.DataFrame, slide_index: int) -> None:

        top = Inches(2.25)
        left = Inches(1.25)
        width = Inches(10.44)
        height = Inches(4.35)
        shapes = self.prs.slides[slide_index].shapes
        rows, cols = data.shape
        table = shapes.add_table(rows + 1, cols, left, top, width,
                                 height).table
        names = data.columns
        for i in range(0, len(names)):
            table.cell(0, i).text = names[i]
        for row in data.itertuples():
            for index in range(data.columns.size):
                value = str(data.iloc[row.Index, index])
                table.cell(int(row.Index + 1), int(index)).text = value

    def save(self):
        self.prs.save(self.name_of_file)


def create_new_presentation(name_new_file: str) -> None:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save(f"./pptx files/{name_new_file}.pptx")

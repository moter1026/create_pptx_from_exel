import os
import pandas as pd

from pptx import Presentation
from pptx.util import Inches


class Present:
    def __init__(self, name_of_file: str, count_slides: int):
        self.name_of_file = name_of_file
        self.count_slides = count_slides
        self.prs = Presentation()

    def add_slide(self, data: dict):
        images = data["images"]
        tables = data["tables"]
        texts = data["texts"]
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])

        for table in tables:
            num_rows, num_columns = table.shape
            slide.shapes.add_table(num_rows, num_columns, Inches(2), Inches(2), Inches(4), Inches(1.5))
# Продолжение следует...




def create_new_presentation(name_new_file: str) -> None:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save(f"./pptx files/{name_new_file}.pptx")
